import hashlib
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.pdfbase import pdfdoc
from reportlab.platypus import Paragraph, Preformatted, SimpleDocTemplate, Spacer

BASE_DIR = Path(__file__).resolve().parents[1]
DOCS_DIR = BASE_DIR / "docs"
SLIDES_DIR = BASE_DIR / "slides"


def _safe_md5(*args, **kwargs):
    return hashlib.md5()


pdfdoc.md5 = _safe_md5


def markdown_to_pdf(source_path: Path, target_path: Path, title: str, compact: bool = False) -> None:
    styles = getSampleStyleSheet()
    body_font_size = 9.2 if compact else 10
    body_leading = 12 if compact else 14
    spacer_height = 4 if compact else 6
    body_style = ParagraphStyle(
        "BodyStyle",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=body_font_size,
        leading=body_leading,
        alignment=TA_LEFT,
        spaceAfter=4 if compact else 6,
    )
    heading1 = ParagraphStyle(
        "Heading1Custom",
        parent=styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=17 if compact else 18,
        leading=20 if compact else 22,
        textColor=colors.HexColor("#1F3A5F"),
        spaceBefore=6 if compact else 8,
        spaceAfter=8 if compact else 10,
    )
    heading2 = ParagraphStyle(
        "Heading2Custom",
        parent=styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=13 if compact else 14,
        leading=16 if compact else 18,
        textColor=colors.HexColor("#2D5D7B"),
        spaceBefore=5 if compact else 6,
        spaceAfter=6 if compact else 8,
    )
    bullet_style = ParagraphStyle(
        "BulletStyle",
        parent=body_style,
        leftIndent=12,
        bulletIndent=0,
    )
    code_style = ParagraphStyle(
        "CodeStyle",
        parent=body_style,
        fontName="Courier",
        fontSize=8.5,
        leading=11,
        backColor=colors.HexColor("#F4F6F8"),
        leftIndent=8,
        rightIndent=8,
    )

    story = [Paragraph(title, heading1), Spacer(1, 6 if compact else 8)]
    in_code_block = False
    code_lines = []

    for raw_line in source_path.read_text().splitlines():
        line = raw_line.rstrip()
        if line.startswith("```"):
            if in_code_block:
                story.append(Preformatted("\n".join(code_lines), code_style))
                story.append(Spacer(1, 8))
                code_lines = []
                in_code_block = False
            else:
                in_code_block = True
            continue

        if in_code_block:
            code_lines.append(line)
            continue

        if not line:
            story.append(Spacer(1, spacer_height))
            continue

        if line.startswith("# "):
            story.append(Paragraph(line[2:], heading1))
        elif line.startswith("## "):
            story.append(Paragraph(line[3:], heading2))
        elif line.startswith("### "):
            story.append(Paragraph(line[4:], heading2))
        elif line.startswith("- "):
            story.append(Paragraph(line[2:], bullet_style, bulletText="•"))
        elif line[0:2].isdigit() and line[1:3] == ". ":
            story.append(Paragraph(line[3:], bullet_style))
        else:
            story.append(Paragraph(line, body_style))

    doc = SimpleDocTemplate(
        str(target_path),
        pagesize=A4,
        topMargin=30 if compact else 36,
        bottomMargin=30 if compact else 36,
        leftMargin=42 if compact else 48,
        rightMargin=42 if compact else 48,
    )
    doc.build(story)


def add_text(slide, left, top, width, height, text, font_size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.text = text
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, bullet_lines, color, font_size=18):
    box = slide.shapes.add_textbox(left, top, width, Inches(3.9))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    first = True
    for line in bullet_lines:
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.bullet = True
        first = False
    return box


def add_card(slide, left, top, width, height, title, body, fill_color, title_color, body_color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.18)
    text_frame.margin_right = Inches(0.18)
    text_frame.margin_top = Inches(0.14)
    text_frame.margin_bottom = Inches(0.12)
    title_paragraph = text_frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.font.size = Pt(16)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = title_color
    body_paragraph = text_frame.add_paragraph()
    body_paragraph.text = body
    body_paragraph.font.size = Pt(12)
    body_paragraph.font.color.rgb = body_color
    return shape


def add_slide_frame(slide, title, page_number, background, title_color, accent):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = background
    add_text(slide, Inches(0.7), Inches(0.45), Inches(9.2), Inches(0.7), title, 28, title_color, bold=True)
    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(1.3), Inches(1.45), Inches(0.12))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()
    add_text(slide, Inches(11.9), Inches(0.42), Inches(0.6), Inches(0.3), f"{page_number:02d}", 10, title_color, align=PP_ALIGN.RIGHT)


def build_presentation(target_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    navy = RGBColor(30, 39, 97)
    mist = RGBColor(247, 244, 234)
    coral = RGBColor(249, 97, 103)
    teal = RGBColor(6, 90, 130)
    sand = RGBColor(231, 232, 209)
    olive = RGBColor(151, 188, 98)
    body = RGBColor(39, 44, 52)
    white = RGBColor(255, 255, 255)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_frame(slide, "BookScope API", 1, navy, white, coral)
    add_text(
        slide,
        Inches(0.75),
        Inches(1.9),
        Inches(7.5),
        Inches(1.4),
        "Google Books ingestion, local catalog ownership, personal bookshelf workflows, explainable recommendations, and reading analytics.",
        24,
        white,
    )
    add_card(slide, Inches(8.7), Inches(1.85), Inches(3.7), Inches(1.25), "Stack", "Django REST Framework\nSQLite\nSimple JWT\nOpenAPI", RGBColor(52, 67, 128), white, white)
    add_text(slide, Inches(0.75), Inches(6.45), Inches(7.8), Inches(0.35), "Repository: github.com/IJF9GAWHY8FGA8/bookscope-api", 11, white)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_frame(slide, "Problem and Scope", 2, mist, body, coral)
    add_card(slide, Inches(0.8), Inches(1.8), Inches(3.7), Inches(2.0), "Why not a proxy?", "Public book APIs expose metadata, but not personal reading state, private bookshelf logic, or coursework-specific analytics.", white, navy, body)
    add_card(slide, Inches(4.8), Inches(1.8), Inches(3.7), Inches(2.0), "Coursework fit", "The project demonstrates CRUD, authentication, data ingestion, testing, API documentation, and version-controlled delivery.", white, navy, body)
    add_card(slide, Inches(8.8), Inches(1.8), Inches(3.7), Inches(2.0), "Chosen scope", "Book, author, genre, bookshelf, review, recommendation, and analytics endpoints on top of local SQLite data.", white, navy, body)
    add_bullets(slide, Inches(0.95), Inches(4.35), Inches(11.2), ["Public catalog browsing", "Authenticated bookshelf and review actions", "Explainable recommendation results", "Analytics for genres, ratings, and reading summary"], body, 18)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_frame(slide, "Architecture Overview", 3, mist, body, teal)
    module_specs = [
        ("accounts", "Register, login, JWT, current user", Inches(0.9)),
        ("catalog", "Books, authors, genres, filters, import", Inches(3.9)),
        ("engagement", "Bookshelf entries and reviews", Inches(6.9)),
        ("analytics", "Recommendations and aggregate endpoints", Inches(9.9)),
    ]
    for name, text, left in module_specs:
        add_card(slide, left, Inches(2.0), Inches(2.35), Inches(1.85), name, text, white, navy, body)
    add_text(slide, Inches(1.05), Inches(4.35), Inches(11.2), Inches(1.2), "Django REST Framework handles routing, serializers, permissions, pagination, filtering, and schema generation. Services isolate recommendation, analytics, and ingestion logic from HTTP views.", 18, body)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_frame(slide, "Data Model", 4, mist, body, coral)
    entity_positions = [
        ("Author", Inches(0.9), Inches(2.0)),
        ("Genre", Inches(3.3), Inches(2.0)),
        ("Book", Inches(5.7), Inches(2.0)),
        ("BookshelfEntry", Inches(8.1), Inches(2.0)),
        ("Review", Inches(10.5), Inches(2.0)),
    ]
    for title, left, top in entity_positions:
        add_card(slide, left, top, Inches(1.8), Inches(1.2), title, "Core entity", sand, navy, body)
    add_text(slide, Inches(0.95), Inches(4.0), Inches(11.3), Inches(1.3), "Many-to-many: Book to Author and Book to Genre. User-linked uniqueness rules: one bookshelf entry per user per book, and one review per user per book.", 18, body)
    add_bullets(slide, Inches(0.95), Inches(5.0), Inches(11.0), ["Google metadata lands in catalog tables", "User behavior lands in bookshelf and review tables", "Analytics derive from both imported metadata and local activity"], body, 16)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_frame(slide, "Google Books Ingestion", 5, mist, body, teal)
    steps = [
        ("1. Fetch", "Read Google Books payloads or local sample files"),
        ("2. Normalize", "Extract titles, authors, categories, ISBNs, dates, and links"),
        ("3. Persist", "Create or update local authors, genres, and books"),
        ("4. Build on local data", "Run CRUD, recommendation, and analytics from SQLite"),
    ]
    for index, (title, text) in enumerate(steps):
        add_card(slide, Inches(0.95 + index * 3.0), Inches(2.1), Inches(2.55), Inches(1.85), title, text, white, navy, body)
    add_text(slide, Inches(0.95), Inches(4.6), Inches(11.2), Inches(1.2), "This keeps the runtime API database-driven and reproducible instead of depending on live external responses during marking.", 18, body)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_frame(slide, "API Documentation and Endpoints", 6, mist, body, coral)
    add_bullets(slide, Inches(0.95), Inches(1.95), Inches(5.9), [
        "Health, schema, and Swagger UI endpoints",
        "Authentication endpoints for register, token, refresh, and current user",
        "Catalog CRUD for books, authors, and genres",
        "Bookshelf and review workflows",
        "Recommendation and analytics endpoints",
    ], body, 17)
    add_card(slide, Inches(7.2), Inches(1.95), Inches(5.05), Inches(2.0), "Submission docs", "OpenAPI YAML\nAPI documentation PDF\nTechnical report PDF\nSlides PPTX", white, navy, body)
    add_card(slide, Inches(7.2), Inches(4.3), Inches(5.05), Inches(1.5), "Permissions", "Public read for catalog and selected analytics. Staff-only catalog writes. Authenticated user ownership for bookshelf and reviews.", sand, navy, body)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_frame(slide, "Recommendation and Analytics", 7, mist, body, teal)
    add_card(slide, Inches(0.95), Inches(1.95), Inches(2.6), Inches(1.65), "Signal 1", "Preferred genres from highly rated books", white, navy, body)
    add_card(slide, Inches(3.75), Inches(1.95), Inches(2.6), Inches(1.65), "Signal 2", "Preferred authors from highly rated books", white, navy, body)
    add_card(slide, Inches(6.55), Inches(1.95), Inches(2.6), Inches(1.65), "Signal 3", "External rating strength and count", white, navy, body)
    add_card(slide, Inches(9.35), Inches(1.95), Inches(2.6), Inches(1.65), "Signal 4", "Local popularity from bookshelf and review activity", white, navy, body)
    add_bullets(slide, Inches(0.95), Inches(4.15), Inches(11.1), [
        "Recommendations exclude books already on the user's shelf or review history",
        "Analytics cover trending books, genre popularity, rating distribution, reading summary, and top authors",
        "The logic is intentionally explainable for easier testing and oral defence",
    ], body, 17)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_frame(slide, "Testing and Error Handling", 8, mist, body, coral)
    add_card(slide, Inches(0.95), Inches(1.95), Inches(5.2), Inches(2.2), "Automated checks", "Registration and auth\nCatalog list and create\nBookshelf uniqueness\nReview uniqueness\nRecommendation exclusion\nAnalytics responses", white, navy, body)
    add_card(slide, Inches(6.55), Inches(1.95), Inches(5.1), Inches(2.2), "Structured errors", '{"error": {"code": "validation_error", "message": "Validation failed.", "details": {"book_id": ["You already have a bookshelf entry for this book."]}}}', sand, navy, body)
    add_text(slide, Inches(0.95), Inches(4.65), Inches(11.0), Inches(1.0), "The API enforces serializer validation, model constraints, permission checks, and consistent HTTP status handling.", 18, body)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_frame(slide, "Version Control Practice", 9, mist, body, teal)
    commits = [
        "Runtime setup and API modules",
        "Google Books sample and seed data",
        "Automated tests",
        "Source documentation and asset builders",
        "Generated PDFs, schema, and presentation deck",
    ]
    for index, label in enumerate(commits):
        top = Inches(1.85 + index * 0.8)
        add_card(slide, Inches(1.15), top, Inches(10.2), Inches(0.58), f"Commit {index + 1}", label, white if index % 2 == 0 else sand, navy, body)
    add_text(slide, Inches(1.15), Inches(6.2), Inches(10.4), Inches(0.45), "The public repository shows staged delivery instead of a single final dump.", 16, body)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_frame(slide, "Deliverables and PythonAnywhere Deployment", 10, mist, body, coral)
    add_card(slide, Inches(0.95), Inches(1.9), Inches(2.7), Inches(1.7), "Deliverables", "Public repo\nREADME\nAPI docs PDF\nTechnical report PDF", white, navy, body)
    add_card(slide, Inches(3.95), Inches(1.9), Inches(2.7), Inches(1.7), "More deliverables", "GenAI appendix\nConversation logs appendix\nSlides PPTX\nTests", white, navy, body)
    add_card(slide, Inches(6.95), Inches(1.9), Inches(2.7), Inches(1.7), "Deployment method", "Manual web app\nVirtualenv\nPlatform WSGI file\nStatic mappings", white, navy, body)
    add_card(slide, Inches(9.95), Inches(1.9), Inches(2.35), Inches(1.7), "Hosting note", "Final publication still needs a PythonAnywhere account and reload step.", sand, navy, body)
    add_text(slide, Inches(0.95), Inches(4.35), Inches(11.1), Inches(1.1), "Target URL: https://YOUR_PYTHONANYWHERE_USERNAME.pythonanywhere.com/", 17, body)

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_frame(slide, "Technical Report Highlights", 11, mist, body, teal)
    add_bullets(slide, Inches(0.95), Inches(1.95), Inches(5.8), [
        "Why Django, DRF, and SQLite were selected",
        "Why Google Books is treated as an external metadata source, not a runtime backend",
        "Why rule-based recommendation is easier to justify and test in this coursework",
        "Where the main limitations and future improvements sit",
    ], body, 17)
    add_card(slide, Inches(7.0), Inches(1.95), Inches(5.2), Inches(2.2), "Submission links inside the report", "Repository link\nAPI docs link\nSlides link\nGenAI appendix link\nConversation logs appendix link", white, navy, body)
    add_text(slide, Inches(0.95), Inches(4.7), Inches(11.0), Inches(1.0), "The report stays within the five-page body limit while still covering architecture, testing, tradeoffs, limitations, and references.", 18, body)

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_frame(slide, "GenAI Usage and Oversight", 12, mist, body, coral)
    add_card(slide, Inches(0.95), Inches(1.95), Inches(5.2), Inches(2.2), "Where AI helped", "Requirement extraction\nPlanning and decomposition\nAPI and documentation structuring\nSuggestion review for tests and materials", white, navy, body)
    add_card(slide, Inches(6.55), Inches(1.95), Inches(5.1), Inches(2.2), "Human control", "Scope was reduced manually\nWeak suggestions were rejected\nTests and exports were verified locally\nFinal tradeoffs were chosen deliberately", sand, navy, body)
    add_text(slide, Inches(0.95), Inches(4.75), Inches(11.1), Inches(0.9), "The repository includes both a GenAI declaration appendix and a representative conversation logs appendix.", 18, body)

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_frame(slide, "Demo Flow, Limitations, and Next Steps", 13, navy, white, coral)
    add_bullets(slide, Inches(0.95), Inches(1.95), Inches(11.0), [
        "Demo flow: register, import data, browse catalog, add bookshelf entry, create review, request recommendations, inspect analytics",
        "Limitations: SQLite scale, heuristic ranking, and a final PythonAnywhere account-side publish step",
        "Next steps: richer ranking, stronger deployment hardening, more advanced filtering and caching",
    ], white, 18)

    prs.save(str(target_path))


def main():
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)

    markdown_to_pdf(
        DOCS_DIR / "api_documentation.md",
        DOCS_DIR / "api_documentation.pdf",
        "BookScope API Documentation",
    )
    markdown_to_pdf(
        DOCS_DIR / "technical_report.md",
        DOCS_DIR / "technical_report.pdf",
        "BookScope API Technical Report",
        compact=True,
    )
    markdown_to_pdf(
        DOCS_DIR / "genai_usage_appendix.md",
        DOCS_DIR / "genai_usage_appendix.pdf",
        "BookScope GenAI Usage Appendix",
    )
    markdown_to_pdf(
        DOCS_DIR / "conversation_logs_appendix.md",
        DOCS_DIR / "conversation_logs_appendix.pdf",
        "BookScope Conversation Logs Appendix",
    )
    build_presentation(SLIDES_DIR / "bookscope_presentation.pptx")


if __name__ == "__main__":
    main()
